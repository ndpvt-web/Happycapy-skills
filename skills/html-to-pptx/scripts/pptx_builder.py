"""
pptx_builder.py — Reusable helpers for building PPTX from extracted HTML design data.

Usage pattern:
    from pptx_builder import Builder
    b = Builder(theme)           # theme = dict of hex colors from HTML CSS vars
    b.add_slide(layout_fn)       # layout_fn receives (slide, b) and places shapes
    b.save("outputs/deck.pptx")

All coordinates use Inches. Slide size: 13.33 × 7.5 in (16:9).
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN


def hex_to_rgb(h: str) -> RGBColor:
    """Convert '#rrggbb' or 'rrggbb' to RGBColor."""
    h = h.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


# ── Default dark-editorial theme (matches the Noir Editorial HTML style) ──────
DEFAULT_THEME = {
    "bg":      "#0b0b0b",
    "bg1":     "#111111",
    "bg2":     "#181818",
    "text":    "#f0ebe2",
    "muted":   "#6b6660",
    "accent":  "#c9a84c",   # gold / primary accent
    "accent2": "#e8c46a",   # lighter accent
    "dim":     "#3a3530",
}


class Builder:
    """Thin wrapper around python-pptx for design-faithful slide creation."""

    W = Inches(13.33)
    H = Inches(7.5)

    def __init__(self, theme: dict = None):
        self.prs = Presentation()
        self.prs.slide_width  = self.W
        self.prs.slide_height = self.H
        self._blank = self.prs.slide_layouts[6]  # truly blank layout

        raw = {**DEFAULT_THEME, **(theme or {})}
        self.c = {k: hex_to_rgb(v) for k, v in raw.items()}

    # ── Slide management ──────────────────────────────────────────────

    def new_slide(self) -> object:
        """Add a blank slide with the theme background filled."""
        slide = self.prs.slides.add_slide(self._blank)
        self._bg(slide)
        return slide

    def save(self, path: str):
        self.prs.save(path)
        print(f"Saved → {path}")

    # ── Core drawing primitives ───────────────────────────────────────

    def _bg(self, slide):
        s = slide.shapes.add_shape(1, 0, 0, self.W, self.H)
        s.fill.solid()
        s.fill.fore_color.rgb = self.c["bg"]
        s.line.fill.background()

    def rect(self, slide, x, y, w, h, fill="bg1", border=None, radius=0):
        """Filled rectangle. fill/border are theme keys or hex strings."""
        s = slide.shapes.add_shape(1, x, y, w, h)
        s.fill.solid()
        s.fill.fore_color.rgb = self._col(fill)
        if border:
            s.line.color.rgb = self._col(border)
            s.line.width = Pt(0.75)
        else:
            s.line.fill.background()
        return s

    def hline(self, slide, x, y, w, color="accent", thickness=Pt(1.5)):
        """Thin horizontal rule (accent line)."""
        s = slide.shapes.add_shape(1, x, y, w, Inches(0.018))
        s.fill.solid()
        s.fill.fore_color.rgb = self._col(color)
        s.line.fill.background()
        return s

    def oval(self, slide, x, y, w, h, fill="bg1"):
        """Oval / circle shape."""
        s = slide.shapes.add_shape(9, x, y, w, h)
        s.fill.solid()
        s.fill.fore_color.rgb = self._col(fill)
        s.line.fill.background()
        return s

    def txt(self, slide, text, x, y, w, h,
            size=14, bold=False, italic=False,
            color="text", align=PP_ALIGN.LEFT, wrap=True):
        """Simple single-run text box."""
        box = slide.shapes.add_textbox(x, y, w, h)
        box.word_wrap = wrap
        tf = box.text_frame
        tf.word_wrap = wrap
        p = tf.paragraphs[0]
        p.alignment = align
        run = p.add_run()
        run.text = text
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.italic = italic
        run.font.color.rgb = self._col(color)
        return box

    def multicolor_txt(self, slide, runs, x, y, w, h,
                       align=PP_ALIGN.LEFT, wrap=True):
        """
        Multi-run text box for mixed colors/styles in one line.
        runs: list of dicts with keys: text, size, bold, italic, color
              e.g. [{"text":"Hello ", "size":48, "bold":True, "color":"text"},
                    {"text":"World",  "size":48, "bold":True, "color":"accent"}]
        """
        box = slide.shapes.add_textbox(x, y, w, h)
        box.word_wrap = wrap
        tf = box.text_frame
        tf.word_wrap = wrap
        p = tf.paragraphs[0]
        p.alignment = align
        for r in runs:
            run = p.add_run()
            run.text = r["text"]
            run.font.size = Pt(r.get("size", 14))
            run.font.bold = r.get("bold", False)
            run.font.italic = r.get("italic", False)
            run.font.color.rgb = self._col(r.get("color", "text"))
        return box

    def slide_number(self, slide, n, total):
        self.txt(slide, f"{n:02d} / {total:02d}",
                 self.W - Inches(1.6), self.H - Inches(0.55),
                 Inches(1.4), Inches(0.38),
                 size=8, color="muted", align=PP_ALIGN.RIGHT)

    def logo(self, slide, name="", x=Inches(0.8), y=Inches(0.38)):
        self.txt(slide, name.upper(), x, y, Inches(3), Inches(0.38),
                 size=8, bold=True, color="muted")

    def eyebrow(self, slide, text, x, y, w=Inches(9)):
        self.txt(slide, text, x, y, w, Inches(0.38),
                 size=9, bold=True, color="accent")

    def card(self, slide, x, y, w, h, fill="bg1", border="dim"):
        return self.rect(slide, x, y, w, h, fill=fill, border=border)

    # ── Layout helpers ────────────────────────────────────────────────

    def grid_cards(self, slide, items, x0, y0, card_w, card_h, gap,
                   cols=None, draw_fn=None):
        """
        Lay out N items in a grid of cards.
        draw_fn(slide, b, item, cx, cy, cw, ch) draws content inside each card.
        """
        n = len(items)
        if cols is None:
            cols = n
        for i, item in enumerate(items):
            col = i % cols
            row = i // cols
            cx = x0 + col * (card_w + gap)
            cy = y0 + row * (card_h + gap)
            self.card(slide, cx, cy, card_w, card_h)
            if draw_fn:
                draw_fn(slide, self, item, cx, cy, card_w, card_h)

    # ── Internal ──────────────────────────────────────────────────────

    def _col(self, key) -> RGBColor:
        """Resolve theme key or raw hex string to RGBColor."""
        if isinstance(key, RGBColor):
            return key
        if key in self.c:
            return self.c[key]
        return hex_to_rgb(key)


# ── HTML theme extractor ──────────────────────────────────────────────────────

import re

def extract_theme_from_html(html: str) -> dict:
    """
    Parse CSS custom properties from an HTML file's :root { } block.
    Returns dict of variable-name → hex-value for color variables.

    Example CSS:   --bg: #0b0b0b;   →   {"bg": "#0b0b0b"}
    Non-color values (clamp(), rem, etc.) are skipped.
    """
    theme = {}
    root_match = re.search(r':root\s*\{([^}]+)\}', html, re.DOTALL)
    if not root_match:
        return theme
    for line in root_match.group(1).splitlines():
        m = re.match(r'\s*--([\w-]+)\s*:\s*(#[0-9a-fA-F]{3,8})\s*;', line)
        if m:
            key = m.group(1).replace("-", "_")
            theme[key] = m.group(2)
    return theme
